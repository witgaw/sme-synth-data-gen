#!/usr/bin/env python3
"""
Generate realistic files from documents.json

Reads the canonical JSON and creates actual .eml, .docx, .xlsx, .pptx, .md, and .pdf files
in the output/ directory. Includes scanned-style PDFs for OCR testing.

Usage:
    uv run generate [--output-dir OUTPUT_DIR]
"""

import argparse
import json
import os
import platform
import random
import subprocess
from datetime import datetime
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
from PIL import Image, ImageFilter
import fitz  # pymupdf


def generate_eml(doc: dict, output_dir: Path) -> Path:
    """Generate RFC 822 .eml file"""
    dt = datetime.fromisoformat(doc["timestamp"])
    date_str = dt.strftime("%a, %d %b %Y %H:%M:%S %z")

    to_list = ", ".join([f"{r['name']} <{r['email']}>" for r in doc["recipients"]])
    cc_list = ", ".join([f"{r['name']} <{r['email']}>" for r in doc.get("cc", [])])

    eml_content = f"""From: {doc["author"]} <{doc["author_email"]}>
To: {to_list}
{"Cc: " + cc_list if cc_list else ""}Date: {date_str}
Subject: {doc["subject"]}
Content-Type: text/plain; charset=utf-8
MIME-Version: 1.0

{doc["body"]}
"""

    filepath = output_dir / doc["filename"]
    filepath.write_text(eml_content, encoding="utf-8")
    return filepath


def generate_md(doc: dict, output_dir: Path) -> Path:
    """Generate markdown file"""
    content = f"# {doc['title']}\n\n"
    content += f"**Author:** {doc['author']}\n"
    content += f"**Date:** {doc['timestamp']}\n\n"

    if "attendees" in doc:
        content += f"**Attendees:** {', '.join(doc['attendees'])}\n"
    if "location" in doc:
        content += f"**Location:** {doc['location']}\n"

    content += "\n---\n\n"
    content += doc.get("content", "")

    if "action_items" in doc:
        content += "\n\n## Action Items\n\n"
        for item in doc["action_items"]:
            content += f"- [ ] {item['task']} ({item['owner']}, due: {item['due']})\n"

    filepath = output_dir / doc["filename"]
    filepath.write_text(content, encoding="utf-8")
    return filepath


def generate_docx(doc: dict, output_dir: Path) -> Path:
    """Generate .docx file"""
    document = Document()
    document.add_heading(doc["title"], 0)

    meta = document.add_paragraph()
    meta.add_run(f"Autor: {doc['author']}\n").italic = True
    dt = datetime.fromisoformat(doc["timestamp"])
    meta.add_run(f"Data: {dt.strftime('%d %B %Y')}").italic = True

    for section in doc.get("sections", []):
        document.add_heading(section["heading"], level=1)
        document.add_paragraph(section["content"])

    filepath = output_dir / doc["filename"]
    document.save(str(filepath))
    return filepath


def generate_xlsx(doc: dict, output_dir: Path) -> Path:
    """Generate .xlsx file"""
    wb = Workbook()
    wb.remove(wb.active)

    for sheet_data in doc.get("sheets", []):
        ws = wb.create_sheet(title=sheet_data["name"][:31])

        # Header row
        for col, header in enumerate(sheet_data["columns"], 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="DDDDDD", end_color="DDDDDD", fill_type="solid")

        # Data rows
        for row_idx, row_data in enumerate(sheet_data["rows"], 2):
            for col_idx, value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=value)

    filepath = output_dir / doc["filename"]
    wb.save(str(filepath))
    return filepath


def generate_pptx(doc: dict, output_dir: Path) -> Path:
    """Generate .pptx file"""
    prs = Presentation()

    for slide_data in doc.get("slides", []):
        if slide_data.get("type") == "title":
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data["title"]
            if "subtitle" in slide_data:
                slide.placeholders[1].text = slide_data["subtitle"]
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data["title"]

            if "bullets" in slide_data:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.text = slide_data["bullets"][0]
                for bullet in slide_data["bullets"][1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
            elif "content" in slide_data:
                body = slide.placeholders[1]
                body.text = slide_data["content"]

    filepath = output_dir / doc["filename"]
    prs.save(str(filepath))
    return filepath


def generate_pdf_easy(doc: dict, output_dir: Path) -> Path:
    """Generate clean, OCR-friendly PDF"""
    filepath = output_dir / doc["filename"]

    # Create PDF with reportlab
    pdf_doc = SimpleDocTemplate(
        str(filepath),
        pagesize=A4,
        rightMargin=2 * cm,
        leftMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )

    styles = getSampleStyleSheet()

    # Custom styles for Polish text
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Heading1"],
        fontSize=16,
        spaceAfter=20,
        fontName="Helvetica-Bold",
    )

    body_style = ParagraphStyle(
        "CustomBody",
        parent=styles["Normal"],
        fontSize=10,
        leading=14,
        fontName="Helvetica",
    )

    story = []

    # Add title
    title = doc.get("title", doc.get("filename", "Dokument"))
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 12))

    # Add content - handle Polish characters by escaping XML entities
    content = doc.get("content", "")
    # Replace problematic characters for XML
    content = content.replace("&", "&amp;")
    content = content.replace("<", "&lt;")
    content = content.replace(">", "&gt;")
    # Convert newlines to HTML breaks for proper rendering
    content = content.replace("\n", "<br/>")

    story.append(Paragraph(content, body_style))

    pdf_doc.build(story)
    return filepath


def apply_scan_effects(img: Image.Image, difficulty: str = "hard") -> Image.Image:
    """Apply scan/degradation effects to an image"""
    # Convert to RGB if necessary
    if img.mode != "RGB":
        img = img.convert("RGB")

    width, height = img.size

    # Rotation (slight for all hard PDFs)
    if difficulty == "hard":
        angle = random.uniform(-2.5, 2.5)
        img = img.rotate(angle, expand=True, fillcolor=(255, 255, 255))

    # Add noise
    if difficulty == "hard":
        import numpy as np

        img_array = np.array(img)
        noise = np.random.normal(0, random.uniform(5, 15), img_array.shape).astype(np.int16)
        img_array = np.clip(img_array.astype(np.int16) + noise, 0, 255).astype(np.uint8)
        img = Image.fromarray(img_array)

    # Slight blur
    if difficulty == "hard":
        blur_radius = random.uniform(0.3, 0.8)
        img = img.filter(ImageFilter.GaussianBlur(radius=blur_radius))

    # Contrast/brightness adjustment
    if difficulty == "hard":
        from PIL import ImageEnhance

        # Reduce contrast slightly
        contrast = ImageEnhance.Contrast(img)
        img = contrast.enhance(random.uniform(0.85, 1.0))

        # Adjust brightness
        brightness = ImageEnhance.Brightness(img)
        img = brightness.enhance(random.uniform(0.9, 1.05))

    return img


def generate_pdf_hard(doc: dict, output_dir: Path) -> Path:
    """Generate scanned-style PDF with degradation effects"""
    # First generate a clean PDF
    temp_pdf_path = output_dir / f"_temp_{doc['filename']}"

    pdf_doc = SimpleDocTemplate(
        str(temp_pdf_path),
        pagesize=A4,
        rightMargin=2 * cm,
        leftMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Heading1"],
        fontSize=14,
        spaceAfter=16,
        fontName="Helvetica-Bold",
    )
    body_style = ParagraphStyle(
        "CustomBody",
        parent=styles["Normal"],
        fontSize=10,
        leading=13,
        fontName="Helvetica",
    )

    story = []
    title = doc.get("title", doc.get("filename", "Dokument"))
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 10))

    content = doc.get("content", "")
    content = content.replace("&", "&amp;")
    content = content.replace("<", "&lt;")
    content = content.replace(">", "&gt;")
    content = content.replace("\n", "<br/>")

    story.append(Paragraph(content, body_style))
    pdf_doc.build(story)

    # Convert PDF to images and apply effects
    pdf_document = fitz.open(str(temp_pdf_path))
    images = []

    # DPI for rendering (lower = more degraded)
    dpi = random.randint(150, 200)
    matrix = fitz.Matrix(dpi / 72, dpi / 72)

    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        pix = page.get_pixmap(matrix=matrix)

        # Convert to PIL Image
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        # Apply scan effects
        img = apply_scan_effects(img, "hard")
        images.append(img)

    pdf_document.close()

    # Save as new PDF from images
    filepath = output_dir / doc["filename"]

    if images:
        # Convert images back to PDF
        first_img = images[0]
        if len(images) > 1:
            first_img.save(
                str(filepath),
                "PDF",
                save_all=True,
                append_images=images[1:],
                resolution=dpi,
                quality=random.randint(70, 85),
            )
        else:
            first_img.save(str(filepath), "PDF", resolution=dpi, quality=random.randint(70, 85))

    # Clean up temp file
    temp_pdf_path.unlink()

    return filepath


def parse_timestamp(ts_str: str) -> datetime:
    """Parse ISO 8601 timestamp string to datetime."""
    # Handle timezone offset format (e.g., +02:00)
    if "+" in ts_str or ts_str.endswith("Z"):
        # Python 3.11+ has fromisoformat support for this, but for 3.10 compatibility:
        ts_str = ts_str.replace("Z", "+00:00")
        # fromisoformat handles +02:00 format in Python 3.11+
        try:
            return datetime.fromisoformat(ts_str)
        except ValueError:
            # Fallback for Python 3.10
            if "+" in ts_str:
                base, tz = ts_str.rsplit("+", 1)
                return datetime.fromisoformat(base)
    return datetime.fromisoformat(ts_str)


def set_file_timestamps(filepath: Path, doc: dict, vary_mtime: bool = True) -> None:
    """
    Set file timestamps to match document metadata.

    Args:
        filepath: Path to the generated file
        doc: Document dict with 'timestamp' field
        vary_mtime: If True, some files get a slightly later mtime (simulating edits)
    """
    ts_str = doc.get("timestamp")
    if not ts_str:
        return

    created_dt = parse_timestamp(ts_str)
    created_ts = created_dt.timestamp()

    # For modification time, sometimes add a small offset to simulate edits
    # ~30% of files have been "edited" after creation
    if vary_mtime and random.random() < 0.3:
        # Add 1-48 hours for edits
        edit_offset = random.randint(3600, 172800)
        mtime_ts = created_ts + edit_offset
    else:
        mtime_ts = created_ts

    # Set atime and mtime (works on all platforms)
    os.utime(filepath, (mtime_ts, mtime_ts))

    # On macOS, also try to set creation time (birthtime)
    if platform.system() == "Darwin":
        try:
            # Format: [[CC]YY]MMDDhhmm[.SS]
            fmt_time = datetime.fromtimestamp(created_ts).strftime("%Y%m%d%H%M.%S")
            # SetFile -d sets creation date (requires Xcode CLI tools)
            subprocess.run(
                ["SetFile", "-d", fmt_time, str(filepath)],
                capture_output=True,
                check=False,
            )
        except FileNotFoundError:
            # SetFile not available, skip creation time
            pass


def main():
    parser = argparse.ArgumentParser(description="Generate files from documents.json")
    parser.add_argument("--output-dir", "-o", default="output", help="Output directory")
    parser.add_argument("--input", "-i", default="dataset/documents.json", help="Input JSON file")
    parser.add_argument(
        "--no-db",
        action="store_true",
        help="Skip SQLite database generation",
    )
    parser.add_argument(
        "--no-timestamps",
        action="store_true",
        help="Skip setting file timestamps to document dates",
    )
    args = parser.parse_args()

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    with open(args.input, "r", encoding="utf-8") as f:
        data = json.load(f)

    docs_to_generate = data["documents"]

    print(f"Generating {len(docs_to_generate)} files to {output_dir}/")

    generated = 0
    skipped = 0

    for doc in docs_to_generate:
        doc_type = doc["type"]
        fmt = doc.get("format", "")

        try:
            if fmt == "pdf":
                # Handle PDF documents
                difficulty = doc.get("pdf_difficulty", "easy")
                if difficulty == "easy":
                    generate_pdf_easy(doc, output_dir)
                else:
                    generate_pdf_hard(doc, output_dir)
            elif fmt == "eml" or "email" in doc_type:
                generate_eml(doc, output_dir)
            elif fmt == "md" or doc_type in ["meeting_notes", "project_kickoff"]:
                generate_md(doc, output_dir)
            elif fmt == "docx" or doc_type in [
                "report_quarterly",
                "report_monthly",
                "report_project",
                "proposal",
            ]:
                generate_docx(doc, output_dir)
            elif fmt == "xlsx" or "spreadsheet" in doc_type:
                generate_xlsx(doc, output_dir)
            elif fmt == "pptx" or "presentation" in doc_type:
                generate_pptx(doc, output_dir)
            else:
                print(f"  Unknown format for {doc['id']}: {doc_type}/{fmt}")
                skipped += 1
                continue

            generated += 1

            # Set file timestamps to match document metadata
            if not args.no_timestamps:
                filepath = output_dir / doc["filename"]
                set_file_timestamps(filepath, doc)

            print(f"  {doc['id']}: {doc['filename']}")

        except Exception as e:
            print(f"  Error generating {doc['id']}: {e}")
            skipped += 1

    print(f"\nDone: {generated} generated, {skipped} skipped")

    # Validation check: ensure we processed all documents
    expected = len(docs_to_generate)
    actual = generated + skipped
    if actual != expected:
        print(f"\nERROR: Expected to process {expected} documents, but processed {actual}")
        raise SystemExit(1)

    if skipped > 0:
        print(f"\nWARNING: {skipped} documents were skipped (missing dependencies?)")

    # Delete old database file before counting (if it exists from a previous run)
    if not args.no_db:
        db_json_path = Path(args.input).parent / "database.json"
        if db_json_path.exists():
            with open(db_json_path, "r", encoding="utf-8") as f:
                db_def = json.load(f)
                db_name = db_def["meta"]["database_name"]
                db_path = output_dir / db_name
                if db_path.exists():
                    db_path.unlink()

    # Count actual files on disk
    actual_files = list(output_dir.glob("*"))
    actual_file_count = len([f for f in actual_files if f.is_file()])

    if actual_file_count != generated:
        print(f"\nERROR: Expected {generated} files on disk, but found {actual_file_count}")
        raise SystemExit(1)

    # Summary
    total_in_json = len(data["documents"])
    if generated == total_in_json:
        print(f"\nValidation passed: all {generated} documents generated")
    else:
        print(f"\nValidation passed: {generated}/{total_in_json} generated ({skipped} skipped)")

    # Generate database by default (unless --no-db)
    if not args.no_db:
        from scripts.generate_database import (
            create_indexes,
            create_schema,
            create_views,
            insert_data,
            verify_database,
        )

        db_json_path = Path(args.input).parent / "database.json"
        with open(db_json_path, "r", encoding="utf-8") as f:
            db_def = json.load(f)

        db_name = db_def["meta"]["database_name"]
        db_path = output_dir / db_name

        print(f"\nGenerating database: {db_path}")

        import sqlite3

        conn = sqlite3.connect(str(db_path))
        try:
            create_schema(conn, db_def["schema"])
            counts = insert_data(conn, db_def["data"])
            create_indexes(conn)
            create_views(conn)

            if verify_database(conn, counts):
                print(f"  Database created: {sum(counts.values())} rows in {len(counts)} tables")
                db_generated = True
            else:
                print("  ERROR: Database verification failed")
                raise SystemExit(1)
        finally:
            conn.close()
    else:
        db_generated = False

    # Final summary
    print("\n" + "=" * 50)
    print("GENERATION SUMMARY")
    print("=" * 50)
    print(f"  Documents: {generated}/{len(docs_to_generate)}")
    print(f"  Database:  {'generated' if db_generated else 'skipped'}")
    print(f"  Timestamps: {'skipped' if args.no_timestamps else 'set to document dates'}")
    print(f"  Output:    {output_dir.absolute()}")
    print("=" * 50)


if __name__ == "__main__":
    main()
